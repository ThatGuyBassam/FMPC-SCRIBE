"""
ingest_slides.py - Ingest PDF or PPTX slide decks into ChromaDB

Modes:
    py -3.11 ingest_slides.py                          -> watch mode (auto)
    py -3.11 ingest_slides.py watch                    -> watch mode (auto)
    py -3.11 ingest_slides.py list                     -> show all ingested files
    py -3.11 ingest_slides.py delete "filename"        -> remove file from DB
    py -3.11 ingest_slides.py file.pdf bacteriologie   -> manual one-shot
    py -3.11 ingest_slides.py file.pptx                -> manual, discipline inferred

Folder structure for watch mode:
    C:\\FMPC_Scribe\\slides\\bacteriologie\\any_name.pdf
    C:\\FMPC_Scribe\\slides\\anatomie\\any_name.pptx
"""

import sys
import os
import time
from datetime import datetime
from collections import Counter

# ── CONFIG ─────────────────────────────────────────────────────────
CHROMA_PATH = r"C:\FMPC_Scribe\chroma_db"
SLIDES_DIR  = r"C:\FMPC_Scribe\slides"
EMBED_MODEL = "paraphrase-multilingual-MiniLM-L12-v2"
KNOWN_DISCIPLINES = [
    "anatomie", "bacteriologie", "histologie", "hematologie",
    "embryologie", "immunologie", "physiologie", "biochimie", "virologie"
]
SUPPORTED_EXT = (".pdf", ".pptx", ".ppt")

# ── DEPENDENCY CHECKS ──────────────────────────────────────────────
try:
    import pdfplumber
except ImportError:
    print("[INGEST] ERROR: pdfplumber not installed.")
    print("  Run: py -3.11 -m pip install pdfplumber")
    sys.exit(1)

try:
    import chromadb
    from chromadb.utils import embedding_functions
except ImportError:
    print("[INGEST] ERROR: chromadb not installed.")
    print("  Run: py -3.11 -m pip install chromadb")
    sys.exit(1)

try:
    from sentence_transformers import SentenceTransformer  # noqa
except ImportError:
    print("[INGEST] ERROR: sentence-transformers not installed.")
    print("  Run: py -3.11 -m pip install sentence-transformers")
    sys.exit(1)

# ── EMBEDDING + CHROMA ─────────────────────────────────────────────
def get_embedding_fn():
    return embedding_functions.SentenceTransformerEmbeddingFunction(
        model_name=EMBED_MODEL
    )

def get_collection(name):
    os.makedirs(CHROMA_PATH, exist_ok=True)
    client = chromadb.PersistentClient(path=CHROMA_PATH)
    col = client.get_or_create_collection(
        name=name,
        embedding_function=get_embedding_fn(),
        metadata={"hnsw:space": "cosine"}
    )
    return col

# ── DISCIPLINE INFERENCE ───────────────────────────────────────────
def infer_discipline(filepath):
    """
    Read discipline from the parent subfolder name.
    slides/bacteriologie/file.pdf -> bacteriologie
    Falls back to filename matching if parent not a known discipline.
    """
    parent = os.path.basename(os.path.dirname(filepath)).lower()
    if parent in KNOWN_DISCIPLINES:
        return parent
    fname = os.path.splitext(os.path.basename(filepath))[0].lower()
    return next((d for d in KNOWN_DISCIPLINES if d in fname), "medecine")

# ── OCR FALLBACK ───────────────────────────────────────────────────
def _ocr_pdf(path):
    """OCR fallback for image-only PDFs using PyMuPDF + Tesseract."""
    try:
        import fitz
        import pytesseract
        from PIL import Image
        import io
    except ImportError as e:
        print(f"[INGEST] OCR dependencies missing: {e}")
        print("  Run: py -3.11 -m pip install pymupdf pytesseract pillow")
        print("  Also install Tesseract: https://github.com/UB-Mannheim/tesseract/wiki")
        return []

    print("[INGEST] Falling back to OCR (image-only PDF)...")
    pages = []
    doc   = fitz.open(path)
    total = len(doc)
    print(f"[INGEST] OCR: {total} pages at 300 DPI...")

    for i, page in enumerate(doc, 1):
        try:
            mat  = fitz.Matrix(300 / 72, 300 / 72)
            pix  = page.get_pixmap(matrix=mat, alpha=False)
            img  = Image.open(io.BytesIO(pix.tobytes("png")))
            text = pytesseract.image_to_string(img, lang="fra", config="--psm 3")
            text = text.strip()
            if len(text) >= 20:
                pages.append((i, text))
                print(f"[INGEST]   Page {i}/{total}: OCR OK ({len(text)} chars)")
            else:
                print(f"[INGEST]   Page {i}/{total}: skipped (too little text)")
        except Exception as e:
            print(f"[INGEST]   Page {i}/{total}: OCR failed — {e}")

    doc.close()
    print(f"[INGEST] OCR complete: {len(pages)} usable pages")
    return pages

# ── EXTRACTION ─────────────────────────────────────────────────────
def extract_pdf(path):
    pages = []
    text_page_count = 0

    with pdfplumber.open(path) as pdf:
        total = len(pdf.pages)
        print(f"[INGEST] PDF: {total} pages")
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if text and len(text.strip()) >= 20:
                pages.append((i, text.strip()))
                text_page_count += 1
            else:
                print(f"[INGEST]   Slide {i}: no embedded text")

    # Fewer than 10% of pages had extractable text — treat as scanned
    if text_page_count < max(1, total * 0.10):
        print(f"[INGEST] Only {text_page_count}/{total} pages had embedded text — treating as scanned PDF")
        pages = _ocr_pdf(path)
    else:
        print(f"[INGEST] Extracted {len(pages)} usable slides")

    return pages

def extract_pptx(path):
    try:
        from pptx import Presentation
    except ImportError:
        print("[INGEST] ERROR: python-pptx not installed.")
        print("  Run: py -3.11 -m pip install python-pptx")
        sys.exit(1)

    pages = []
    prs = Presentation(path)
    total = len(prs.slides)
    print(f"[INGEST] PPTX: {total} slides")

    for i, slide in enumerate(prs.slides, 1):
        lines = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = " ".join(run.text for run in para.runs).strip()
                if line:
                    lines.append(line)
        text = "\n".join(lines).strip()
        if len(text) >= 20:
            pages.append((i, text))
        else:
            print(f"[INGEST]   Slide {i}: skipped (too short)")

    print(f"[INGEST] Extracted {len(pages)} usable slides")
    return pages

def extract(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return extract_pdf(path)
    elif ext in (".pptx", ".ppt"):
        return extract_pptx(path)
    else:
        print(f"[INGEST] ERROR: Unsupported format: {ext}")
        sys.exit(1)

# ── INGEST FILE ────────────────────────────────────────────────────
def ingest_file(path, discipline):
    filename     = os.path.basename(path)
    name_noext   = os.path.splitext(filename)[0]
    today        = datetime.now().strftime("%Y-%m-%d")
    ext          = os.path.splitext(path)[1].lower()
    filetype     = "PPTX" if ext in (".pptx", ".ppt") else "PDF"

    print(f"[INGEST] --- Starting ingestion ---")
    print(f"[INGEST] File:       {filename} ({filetype})")
    print(f"[INGEST] Discipline: {discipline}")
    print(f"[INGEST] Date:       {today}")

    pages = extract(path)
    if not pages:
        print(f"[INGEST] ERROR: No text extracted. File may be image-only.")
        return

    col = get_collection("slides")

    documents, metadatas, ids = [], [], []
    for page_num, text in pages:
        doc_id = f"{name_noext}_p{page_num}"
        existing = col.get(ids=[doc_id])
        if existing and existing["ids"]:
            print(f"[INGEST]   Slide {page_num}: already in DB, skipping.")
            continue
        documents.append(text)
        metadatas.append({
            "source":     "slide",
            "discipline": discipline,
            "filename":   name_noext,
            "date":       today,
            "page":       page_num,
        })
        ids.append(doc_id)

    if not documents:
        print("[INGEST] All slides already ingested.")
        return

    BATCH = 50
    total = len(documents)
    print(f"[INGEST] Ingesting {total} slides...")
    for i in range(0, total, BATCH):
        col.add(
            documents=documents[i:i+BATCH],
            metadatas=metadatas[i:i+BATCH],
            ids=ids[i:i+BATCH]
        )
        print(f"[INGEST]   {min(i+BATCH, total)}/{total} done")

    print(f"[INGEST] --- DONE: {total} slides added. Collection: {col.count()} total ---")

# ── TRANSCRIPT INGEST (called from processor.py) ───────────────────
def ingest_transcript(txt_path, discipline, label, professor=None):
    print(f"[INGEST] Ingesting transcript: {os.path.basename(txt_path)}")

    with open(txt_path, "r", encoding="utf-8") as f:
        raw = f.read()

    if "=" * 10 in raw:
        raw = raw.split("=" * 10, 1)[-1].strip()

    paragraphs = [p.strip() for p in raw.split("\n\n") if len(p.strip()) >= 30]
    if not paragraphs:
        print("[INGEST] No usable paragraphs found.")
        return

    col          = get_collection("transcripts")
    today        = datetime.now().strftime("%Y-%m-%d")
    name_noext   = os.path.splitext(os.path.basename(txt_path))[0]

    documents, metadatas, ids = [], [], []
    for i, para in enumerate(paragraphs):
        doc_id = f"{name_noext}_chunk{i}"
        existing = col.get(ids=[doc_id])
        if existing and existing["ids"]:
            continue
        documents.append(para)
        metadatas.append({
            "source":      "transcript",
            "discipline":  discipline,
            "filename":    name_noext,
            "label":       label,
            "professor":   professor or "unknown",
            "date":        today,
            "chunk_index": i,
        })
        ids.append(doc_id)

    if not documents:
        print("[INGEST] Transcript already ingested.")
        return

    BATCH = 50
    total = len(documents)
    for i in range(0, total, BATCH):
        col.add(
            documents=documents[i:i+BATCH],
            metadatas=metadatas[i:i+BATCH],
            ids=ids[i:i+BATCH]
        )
    print(f"[INGEST] Transcript ingested: {total} chunks. Collection: {col.count()} total")

# ── LIST ───────────────────────────────────────────────────────────
def cmd_list():
    os.makedirs(CHROMA_PATH, exist_ok=True)
    client = chromadb.PersistentClient(path=CHROMA_PATH)
    ef = get_embedding_fn()
    for col_name in ["slides", "transcripts"]:
        try:
            col     = client.get_collection(name=col_name, embedding_function=ef)
            results = col.get(include=["metadatas"])
            if not results["ids"]:
                print(f"\n[{col_name.upper()}] Empty.")
                continue
            files = Counter(m["filename"] for m in results["metadatas"])
            print(f"\n[{col_name.upper()}] {col.count()} chunks across {len(files)} files:")
            for fname, count in sorted(files.items()):
                disc = next(
                    (m["discipline"] for m in results["metadatas"] if m["filename"] == fname),
                    "?"
                )
                print(f"   {disc:15} {fname}  ({count} chunks)")
        except Exception:
            print(f"\n[{col_name.upper()}] Not created yet.")

# ── DELETE ─────────────────────────────────────────────────────────
def cmd_delete(target):
    os.makedirs(CHROMA_PATH, exist_ok=True)
    client  = chromadb.PersistentClient(path=CHROMA_PATH)
    ef      = get_embedding_fn()
    deleted = 0
    for col_name in ["slides", "transcripts"]:
        try:
            col     = client.get_collection(name=col_name, embedding_function=ef)
            results = col.get(include=["metadatas"])
            to_del  = [
                rid for rid, meta in zip(results["ids"], results["metadatas"])
                if meta.get("filename") == target
            ]
            if to_del:
                col.delete(ids=to_del)
                print(f"[DELETE] Removed {len(to_del)} chunks from '{col_name}': {target}")
                deleted += len(to_del)
            else:
                print(f"[DELETE] Nothing in '{col_name}' for: {target}")
        except Exception as e:
            print(f"[DELETE] Could not access '{col_name}': {e}")
    if deleted:
        print(f"[DELETE] Done. {deleted} total chunks removed.")
    else:
        print("[DELETE] Nothing deleted. Check name with: py -3.11 ingest_slides.py list")

# ── WATCHER ────────────────────────────────────────────────────────
def cmd_watch():
    try:
        from watchdog.observers import Observer
        from watchdog.events import FileSystemEventHandler
    except ImportError:
        print("[WATCHER] ERROR: watchdog not installed.")
        print("  Run: py -3.11 -m pip install watchdog")
        sys.exit(1)

    os.makedirs(SLIDES_DIR, exist_ok=True)
    recently_processed = set()

    class SlideHandler(FileSystemEventHandler):
        def on_created(self, event):
            if event.is_directory:
                return
            path = event.src_path
            if not any(path.lower().endswith(e) for e in SUPPORTED_EXT):
                return
            if path in recently_processed:
                return

            time.sleep(2)
            try:
                s1 = os.path.getsize(path)
                time.sleep(1)
                s2 = os.path.getsize(path)
                if s1 != s2:
                    time.sleep(3)
            except Exception:
                return

            recently_processed.add(path)
            discipline = infer_discipline(path)
            print(f"\n[WATCHER] New file: {os.path.basename(path)}")
            print(f"[WATCHER] Discipline: {discipline} (from subfolder)")
            try:
                ingest_file(path, discipline)
            except Exception as e:
                print(f"[WATCHER] ERROR: {e}")

    print(f"[WATCHER] --- Slide Watcher Active ---")
    print(f"[WATCHER] Watching: {SLIDES_DIR}")
    print(f"[WATCHER] Put PDFs/PPTXs in subfolders: slides/bacteriologie/file.pdf")
    print(f"[WATCHER] Press Ctrl+C to stop.\n")

    # Startup scan — ingest anything already sitting in the folder
    print("[WATCHER] Scanning for un-ingested files...")
    found = False
    for root, dirs, files in os.walk(SLIDES_DIR):
        for fname in files:
            if not any(fname.lower().endswith(e) for e in SUPPORTED_EXT):
                continue
            fpath      = os.path.join(root, fname)
            discipline = infer_discipline(fpath)
            name_noext = os.path.splitext(fname)[0]
            try:
                col      = get_collection("slides")
                existing = col.get(ids=[f"{name_noext}_p1"])
                if existing and existing["ids"]:
                    print(f"[WATCHER]   Already ingested: {fname}")
                    continue
            except Exception:
                pass
            subfolder = os.path.basename(root)
            print(f"[WATCHER]   Ingesting: {subfolder}/{fname} ({discipline})")
            found = True
            try:
                ingest_file(fpath, discipline)
            except Exception as e:
                print(f"[WATCHER]   ERROR: {e}")

    if not found:
        print("[WATCHER] No new files found. Waiting...\n")

    observer = Observer()
    observer.schedule(SlideHandler(), path=SLIDES_DIR, recursive=True)
    observer.start()
    try:
        while True:
            time.sleep(1)
    except KeyboardInterrupt:
        observer.stop()
        print("\n[WATCHER] Stopped.")
    observer.join()

# ── ENTRY POINT ────────────────────────────────────────────────────
if __name__ == "__main__":

    # No args or "watch" -> watch mode
    if len(sys.argv) == 1 or sys.argv[1] == "watch":
        cmd_watch()
        sys.exit(0)

    # list
    if sys.argv[1] == "list":
        cmd_list()
        sys.exit(0)

    # delete <filename>
    if sys.argv[1] == "delete":
        if len(sys.argv) < 3:
            print("Usage: py -3.11 ingest_slides.py delete \"filename_without_extension\"")
            sys.exit(1)
        cmd_delete(sys.argv[2])
        sys.exit(0)

    # manual ingest <path> [discipline]
    file_path = sys.argv[1]
    if not os.path.exists(file_path):
        print(f"[INGEST] ERROR: File not found: {file_path}")
        sys.exit(1)
    if not any(file_path.lower().endswith(e) for e in SUPPORTED_EXT):
        print(f"[INGEST] ERROR: Unsupported format. Use .pdf or .pptx")
        sys.exit(1)

    if len(sys.argv) >= 3:
        discipline = sys.argv[2].lower().strip()
    else:
        discipline = infer_discipline(file_path)
        print(f"[INGEST] Discipline inferred: {discipline}")

    ingest_file(file_path, discipline)print(f"[WATCHER]   Already ingested: {fname}")
                    continue
            except Exception:
                pass
            subfolder = os.path.basename(root)
            print(f"[WATCHER]   Ingesting: {subfolder}/{fname} ({discipline})")
            found = True
            try:
                ingest_file(fpath, discipline)
            except Exception as e:
                print(f"[WATCHER]   ERROR: {e}")

    if not found:
        print("[WATCHER] No new files found. Waiting...\n")

    observer = Observer()
    observer.schedule(SlideHandler(), path=SLIDES_DIR, recursive=True)
    observer.start()
    try:
        while True:
            time.sleep(1)
    except KeyboardInterrupt:
        observer.stop()
        print("\n[WATCHER] Stopped.")
    observer.join()

# ── ENTRY POINT ────────────────────────────────────────────────────
if __name__ == "__main__":

    # No args or "watch" -> watch mode
    if len(sys.argv) == 1 or sys.argv[1] == "watch":
        cmd_watch()
        sys.exit(0)

    # list
    if sys.argv[1] == "list":
        cmd_list()
        sys.exit(0)

    # delete <filename>
    if sys.argv[1] == "delete":
        if len(sys.argv) < 3:
            print("Usage: py -3.11 ingest_slides.py delete \"filename_without_extension\"")
            sys.exit(1)
        cmd_delete(sys.argv[2])
        sys.exit(0)

    # manual ingest <path> [discipline]
    file_path = sys.argv[1]
    if not os.path.exists(file_path):
        print(f"[INGEST] ERROR: File not found: {file_path}")
        sys.exit(1)
    if not any(file_path.lower().endswith(e) for e in SUPPORTED_EXT):
        print(f"[INGEST] ERROR: Unsupported format. Use .pdf or .pptx")
        sys.exit(1)

    if len(sys.argv) >= 3:
        discipline = sys.argv[2].lower().strip()
    else:
        discipline = infer_discipline(file_path)
        print(f"[INGEST] Discipline inferred: {discipline}")

    ingest_file(file_path, discipline)
