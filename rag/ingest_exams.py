"""
ingest_exams.py - Ingest past exam PDFs into ChromaDB

Usage:
    py -3.11 ingest_exams.py                                  -> watch mode
    py -3.11 ingest_exams.py list                             -> show all ingested exams
    py -3.11 ingest_exams.py delete "filename"                -> remove from DB
    py -3.11 ingest_exams.py file.pdf bacteriologie 2024      -> manual, with year
    py -3.11 ingest_exams.py file.pdf bacteriologie           -> manual, year inferred

Folder structure for watch mode:
    C:\\FMPC_Scribe\\exams\\bacteriologie\\examen_2024.pdf
    C:\\FMPC_Scribe\\exams\\anatomie\\controle_2023.pdf
"""

import sys
import os
import time
from datetime import datetime
from collections import Counter

# ── CONFIG ─────────────────────────────────────────────────────────
CHROMA_PATH = r"C:\FMPC_Scribe\chroma_db"
EXAMS_DIR   = r"C:\FMPC_Scribe\exams"
EMBED_MODEL = "paraphrase-multilingual-MiniLM-L12-v2"
KNOWN_DISCIPLINES = [
    "anatomie", "bacteriologie", "histologie", "hematologie",
    "embryologie", "immunologie", "physiologie", "biochimie", "virologie"
]

# Exam block folder names map to individual disciplines
EXAM_BLOCK_MAP = {
    "histologie-embryologie":    ["histologie", "embryologie"],
    "hematologie-immunologie":   ["hematologie", "immunologie"],
    "microbiologie":             ["bacteriologie", "virologie"],
    "anatomie":                  ["anatomie"],
    "physiologie":               ["physiologie"],
}
SUPPORTED_EXT = (".pdf", ".pptx", ".ppt")

# ── DEPENDENCY CHECKS ──────────────────────────────────────────────
try:
    import pdfplumber
except ImportError:
    print("[EXAMS] ERROR: pdfplumber not installed.")
    print("  Run: py -3.11 -m pip install pdfplumber")
    sys.exit(1)

try:
    import chromadb
    from chromadb.utils import embedding_functions
except ImportError:
    print("[EXAMS] ERROR: chromadb not installed.")
    sys.exit(1)

try:
    from sentence_transformers import SentenceTransformer  # noqa
except ImportError:
    print("[EXAMS] ERROR: sentence-transformers not installed.")
    sys.exit(1)

# ── EMBEDDING + CHROMA ─────────────────────────────────────────────
def get_collection():
    os.makedirs(CHROMA_PATH, exist_ok=True)
    client = chromadb.PersistentClient(path=CHROMA_PATH)
    ef = embedding_functions.SentenceTransformerEmbeddingFunction(model_name=EMBED_MODEL)
    return client.get_or_create_collection(
        name="exams",
        embedding_function=ef,
        metadata={"hnsw:space": "cosine"}
    )

# ── DISCIPLINE + YEAR INFERENCE ────────────────────────────────────
def infer_discipline(filepath):
    parent = os.path.basename(os.path.dirname(filepath)).lower()
    if parent in KNOWN_DISCIPLINES:
        return parent
    fname = os.path.splitext(os.path.basename(filepath))[0].lower()
    return next((d for d in KNOWN_DISCIPLINES if d in fname), "medecine")

def infer_year(filepath):
    import re
    fname = os.path.basename(filepath)
    match = re.search(r"(20\d{2})", fname)
    return match.group(1) if match else "unknown"

# ── EXTRACTION ─────────────────────────────────────────────────────
def extract_pdf(path):
    pages = []
    with pdfplumber.open(path) as pdf:
        total = len(pdf.pages)
        print(f"[EXAMS] PDF: {total} pages")
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if text:
                text = text.strip()
                if len(text) >= 20:
                    pages.append((i, text))
                else:
                    print(f"[EXAMS]   Page {i}: skipped (too short)")
            else:
                print(f"[EXAMS]   Page {i}: skipped (no text)")
    print(f"[EXAMS] Extracted {len(pages)} usable pages")
    return pages

def extract_pptx(path):
    try:
        from pptx import Presentation
    except ImportError:
        print("[EXAMS] ERROR: python-pptx not installed.")
        sys.exit(1)
    pages = []
    prs = Presentation(path)
    total = len(prs.slides)
    print(f"[EXAMS] PPTX: {total} slides")
    for i, slide in enumerate(prs.slides, 1):
        lines = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = " ".join(run.text for run in para.runs).strip()
                if line:
                    lines.append(line)
        text = "\n".join(lines).strip()
        if len(text) >= 20:
            pages.append((i, text))
    print(f"[EXAMS] Extracted {len(pages)} usable slides")
    return pages

def extract(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return extract_pdf(path)
    elif ext in (".pptx", ".ppt"):
        return extract_pptx(path)
    else:
        print(f"[EXAMS] ERROR: Unsupported format: {ext}")
        sys.exit(1)

# ── INGEST ─────────────────────────────────────────────────────────
def ingest_exam(path, discipline, year):
    filename   = os.path.basename(path)
    name_noext = os.path.splitext(filename)[0]
    today      = datetime.now().strftime("%Y-%m-%d")

    print(f"[EXAMS] --- Starting ingestion ---")
    print(f"[EXAMS] File:       {filename}")
    print(f"[EXAMS] Discipline: {discipline}")
    print(f"[EXAMS] Year:       {year}")

    pages = extract(path)
    if not pages:
        print("[EXAMS] ERROR: No text extracted.")
        return

    col = get_collection()
    documents, metadatas, ids = [], [], []

    for page_num, text in pages:
        doc_id = f"{name_noext}_p{page_num}"
        existing = col.get(ids=[doc_id])
        if existing and existing["ids"]:
            print(f"[EXAMS]   Page {page_num}: already in DB, skipping.")
            continue
        documents.append(text)
        metadatas.append({
            "source":     "exam",
            "discipline": discipline,
            "filename":   name_noext,
            "year":       year,
            "date":       today,
            "page":       page_num,
        })
        ids.append(doc_id)

    if not documents:
        print("[EXAMS] All pages already ingested.")
        return

    BATCH = 50
    total = len(documents)
    print(f"[EXAMS] Ingesting {total} pages...")
    for i in range(0, total, BATCH):
        col.add(
            documents=documents[i:i+BATCH],
            metadatas=metadatas[i:i+BATCH],
            ids=ids[i:i+BATCH]
        )
        print(f"[EXAMS]   {min(i+BATCH, total)}/{total} done")

    print(f"[EXAMS] --- DONE: {total} pages added. Collection: {col.count()} total ---")

# ── LIST ───────────────────────────────────────────────────────────
def cmd_list():
    try:
        col     = get_collection()
        results = col.get(include=["metadatas"])
        if not results["ids"]:
            print("[EXAMS] Empty.")
            return
        files = Counter(m["filename"] for m in results["metadatas"])
        print(f"\n[EXAMS] {col.count()} chunks across {len(files)} files:")
        for fname, count in sorted(files.items()):
            meta = next(m for m in results["metadatas"] if m["filename"] == fname)
            print(f"   {meta.get('discipline', '?'):15} {meta.get('year', '?'):6} {fname}  ({count} chunks)")
    except Exception as e:
        print(f"[EXAMS] Not created yet or error: {e}")

# ── DELETE ─────────────────────────────────────────────────────────
def cmd_delete(target):
    try:
        col     = get_collection()
        results = col.get(include=["metadatas"])
        to_del  = [
            rid for rid, meta in zip(results["ids"], results["metadatas"])
            if meta.get("filename") == target
        ]
        if to_del:
            col.delete(ids=to_del)
            print(f"[DELETE] Removed {len(to_del)} chunks: {target}")
        else:
            print(f"[DELETE] Nothing found for: {target}")
            print("[DELETE] Check with: py -3.11 ingest_exams.py list")
    except Exception as e:
        print(f"[DELETE] Error: {e}")

# ── WATCHER ────────────────────────────────────────────────────────
def cmd_watch():
    try:
        from watchdog.observers import Observer
        from watchdog.events import FileSystemEventHandler
    except ImportError:
        print("[WATCHER] ERROR: watchdog not installed.")
        sys.exit(1)

    os.makedirs(EXAMS_DIR, exist_ok=True)
    recently_processed = set()

    class ExamHandler(FileSystemEventHandler):
        def on_created(self, event):
            if event.is_directory:
                return
            path = event.src_path
            if not any(path.lower().endswith(e) for e in SUPPORTED_EXT):
                return
            if path in recently_processed:
                return
            time.sleep(2)
            try:
                s1 = os.path.getsize(path)
                time.sleep(1)
                s2 = os.path.getsize(path)
                if s1 != s2:
                    time.sleep(3)
            except Exception:
                return
            recently_processed.add(path)
            discipline = infer_discipline(path)
            year       = infer_year(path)
            print(f"\n[WATCHER] New exam: {os.path.basename(path)}")
            print(f"[WATCHER] Discipline: {discipline} | Year: {year}")
            try:
                ingest_exam(path, discipline, year)
            except Exception as e:
                print(f"[WATCHER] ERROR: {e}")

    print(f"[WATCHER] --- Exam Watcher Active ---")
    print(f"[WATCHER] Watching: {EXAMS_DIR}")
    print(f"[WATCHER] Put exams in subfolders: exams/bacteriologie/examen_2024.pdf")
    print(f"[WATCHER] Press Ctrl+C to stop.\n")

    # Startup scan
    print("[WATCHER] Scanning for un-ingested exams...")
    for root, dirs, files in os.walk(EXAMS_DIR):
        for fname in files:
            if not any(fname.lower().endswith(e) for e in SUPPORTED_EXT):
                continue
            fpath      = os.path.join(root, fname)
            discipline = infer_discipline(fpath)
            year       = infer_year(fpath)
            name_noext = os.path.splitext(fname)[0]
            try:
                col      = get_collection()
                existing = col.get(ids=[f"{name_noext}_p1"])
                if existing and existing["ids"]:
                    print(f"[WATCHER]   Already ingested: {fname}")
                    continue
            except Exception:
                pass
            print(f"[WATCHER]   Ingesting: {fname} ({discipline}, {year})")
            try:
                ingest_exam(fpath, discipline, year)
            except Exception as e:
                print(f"[WATCHER]   ERROR: {e}")

    observer = Observer()
    observer.schedule(ExamHandler(), path=EXAMS_DIR, recursive=True)
    observer.start()
    try:
        while True:
            time.sleep(1)
    except KeyboardInterrupt:
        observer.stop()
        print("\n[WATCHER] Stopped.")
    observer.join()

# ── ENTRY POINT ────────────────────────────────────────────────────
if __name__ == "__main__":

    if len(sys.argv) == 1 or sys.argv[1] == "watch":
        cmd_watch()
        sys.exit(0)

    if sys.argv[1] == "list":
        cmd_list()
        sys.exit(0)

    if sys.argv[1] == "delete":
        if len(sys.argv) < 3:
            print("Usage: py -3.11 ingest_exams.py delete \"filename_without_extension\"")
            sys.exit(1)
        cmd_delete(sys.argv[2])
        sys.exit(0)

    file_path = sys.argv[1]
    if not os.path.exists(file_path):
        print(f"[EXAMS] ERROR: File not found: {file_path}")
        sys.exit(1)
    if not any(file_path.lower().endswith(e) for e in SUPPORTED_EXT):
        print(f"[EXAMS] ERROR: Unsupported format. Use .pdf or .pptx")
        sys.exit(1)

    discipline = sys.argv[2].lower().strip() if len(sys.argv) >= 3 else infer_discipline(file_path)
    year       = sys.argv[3].strip() if len(sys.argv) >= 4 else infer_year(file_path)

    print(f"[EXAMS] Discipline: {discipline} | Year: {year}")
    ingest_exam(file_path, discipline, year)
